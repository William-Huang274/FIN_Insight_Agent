"""Portable report delivery using mature Markdown, plotting and Office libraries.

Input is an already submitted report, never arbitrary code/URLs/files. Exports
do not invoke a model, mutate the report, or imply human approval.
"""
from __future__ import annotations

import io
from pathlib import Path
import re
from xml.sax.saxutils import escape


def readable_report(report):
    text = report["narrative_markdown"]
    references = []
    for number, (key, citation) in enumerate(report.get("citations", {}).items(), 1):
        text = text.replace("[" + key + "]", f"[{number}]")
        titles, urls = [], []
        for source in citation.get("sources", []):
            if source.get("title"):
                titles.append(source["title"])
            urls.extend(([source["source_url"]] if source.get("source_url") else []) + list(source.get("citation_urls") or []))
        references.append(f"[{number}] " + "；".join(dict.fromkeys(titles)) + "\n" + "\n".join(dict.fromkeys(urls)))
    def source_urls(value):
        if isinstance(value, dict):
            for key, item in value.items():
                if key in {"source_url", "url"} and isinstance(item, str) and item.startswith(("https://", "http://")):
                    yield item
                elif key == "citation_urls" and isinstance(item, list):
                    yield from (url for url in item if isinstance(url, str) and url.startswith(("https://", "http://")))
                else:
                    yield from source_urls(item)
        elif isinstance(value, list):
            for item in value:
                yield from source_urls(item)
    for chart in report.get("charts", []):
        for point in chart["points"]:
            references.append(f"图表《{chart['title']}》· {point['label']} / {point['series']}：{point['source_id']}\n"
                + "\n".join(dict.fromkeys(source_urls(point["provenance"]))))
    return text, references


def markdown_blocks(text):
    from markdown_it import MarkdownIt
    tokens = MarkdownIt("commonmark").enable("table").parse(text)
    blocks, index = [], 0
    def plain(token):
        if not token.children:
            return token.content
        return "".join("\n" if t.type in {"softbreak", "hardbreak"} else t.content
            for t in token.children if t.type in {"text", "code_inline", "softbreak", "hardbreak"})
    while index < len(tokens):
        token = tokens[index]
        if token.type == "table_open":
            rows = []
            while index < len(tokens) and tokens[index].type != "table_close":
                if tokens[index].type == "tr_open":
                    rows.append([])
                elif tokens[index].type == "inline" and rows:
                    rows[-1].append(plain(tokens[index]))
                index += 1
            if rows:
                blocks.append(("table", rows))
        elif token.type == "heading_open" and index + 1 < len(tokens):
            blocks.append((token.tag, plain(tokens[index + 1])))
            index += 1
        elif token.type == "inline":
            blocks.append(("p", plain(token)))
        elif token.type in {"fence", "code_block"}:
            blocks.append(("p", token.content))
        index += 1
    return blocks


def chart_png(chart):
    from matplotlib.figure import Figure
    from matplotlib.backends.backend_agg import FigureCanvasAgg
    from matplotlib.font_manager import FontProperties
    font_path = next((str(p) for p in (Path("C:/Windows/Fonts/msyh.ttc"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc")) if p.is_file()), None)
    font = FontProperties(fname=font_path) if font_path else FontProperties()
    figure = Figure(figsize=(10, 4.8), dpi=140, facecolor="white", layout="constrained")
    axis = figure.subplots()
    labels = list(dict.fromkeys(p["label"] for p in chart["points"]))
    series = list(dict.fromkeys(p["series"] for p in chart["points"]))
    colors = ["#187b80", "#1c365b", "#c27d3e", "#67788c"]
    for group, name in enumerate(series):
        data = {p["label"]: p["value"] for p in chart["points"] if p["series"] == name}
        indices = [i for i, label in enumerate(labels) if label in data]
        values = [data[labels[i]] for i in indices]
        if chart["kind"] == "line":
            axis.plot(indices, values, marker="o", linewidth=2.6, label=name, color=colors[group % len(colors)])
        else:
            width = .72 / len(series)
            bars = axis.bar([i - .36 + width * (group + .5) for i in indices], values, width=width,
                label=name, color=colors[group % len(colors)])
            axis.bar_label(bars, fmt="%.2f", padding=4, fontsize=10, fontproperties=font)
    axis.set_xticks(range(len(labels)), labels, fontproperties=font, rotation=15 if len(labels) > 5 else 0)
    axis.set_ylabel(chart["unit"], fontproperties=font)
    axis.set_title(chart["title"], loc="left", fontproperties=font, fontsize=16, pad=20)
    axis.spines[["top", "right"]].set_visible(False)
    axis.set_axisbelow(True)
    axis.grid(axis="y", color="#e6ebef", linewidth=.8)
    axis.margins(y=.2)
    if len(series) > 1 or series[0]:
        axis.legend(prop=font, frameon=False)
    output = io.BytesIO()
    FigureCanvasAgg(figure).print_png(output)
    return output.getvalue()


def export_report(report, format, *, review_status="待人工审阅"):
    text, references = readable_report(report)
    charts = report.get("charts", [])
    blocks = markdown_blocks(text)
    title = report["title"]
    if format == "md":
        body = f"# {title}\n\n{review_status}\n\n{text}"
        for chart in charts:
            body += f"\n\n## {chart['title']}\n\n{chart['interpretation']}\n\n单位：{chart['unit']}\n\n| 项目 | 系列 | 数值 |\n|---|---|---:|\n"
            body += "\n".join(f"| {p['label']} | {p['series']} | {p['value']:g} |" for p in chart["points"])
        body += "\n\n## 来源\n\n" + "\n\n".join(references)
        return body.encode("utf-8"), "text/markdown; charset=utf-8"
    if format == "pdf":
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
        font_path = Path("C:/Windows/Fonts/msyh.ttc")
        font = "FinCJK"
        if font not in pdfmetrics.getRegisteredFontNames():
            if font_path.is_file():
                pdfmetrics.registerFont(TTFont(font, str(font_path), subfontIndex=0))
            else:
                font = "STSong-Light"
                pdfmetrics.registerFont(UnicodeCIDFont(font))
        normal = ParagraphStyle("Body", fontName=font, fontSize=10, leading=17, spaceAfter=9, wordWrap="CJK")
        small = ParagraphStyle("Source", parent=normal, fontSize=8, leading=12, splitLongWords=True)
        heading = ParagraphStyle("Heading", parent=normal, fontSize=15, leading=22, spaceBefore=17, keepWithNext=True)
        title_style = ParagraphStyle("Title", parent=normal, fontSize=22, leading=31, spaceAfter=14)
        story = [Paragraph(escape(title), title_style), Paragraph(escape(review_status), small), Spacer(1, 12)]
        def para(value, style=normal):
            return Paragraph(escape(str(value)).replace("\n", "<br/>"), style)
        for kind, value in blocks:
            if kind == "table":
                width = max(len(row) for row in value)
                rows = [[para(cell, small) for cell in row] + [""]*(width-len(row)) for row in value]
                table = Table(rows, colWidths=[470/width]*width, repeatRows=1, hAlign="LEFT")
                table.setStyle(TableStyle([("BACKGROUND", (0,0), (-1,0), colors.HexColor("#e9eff4")),
                    ("GRID", (0,0), (-1,-1), .35, colors.HexColor("#d9d9d9")), ("VALIGN", (0,0), (-1,-1), "MIDDLE"),
                    ("TOPPADDING", (0,0), (-1,-1), 7), ("BOTTOMPADDING", (0,0), (-1,-1), 7)]))
                story.extend([table, Spacer(1, 12)])
            else:
                story.append(para(value, heading if kind.startswith("h") else normal))
        if charts:
            story.append(para("图表与出处", heading))
        for chart in charts:
            story.extend([Image(io.BytesIO(chart_png(chart)), width=470, height=225.6), para(chart["interpretation"], small)])
        story.append(para("来源", heading))
        story.extend(para(ref, small) for ref in references)
        output = io.BytesIO()
        def footer(canvas, doc):
            canvas.setFont(font, 8)
            canvas.setFillColor(colors.HexColor("#657386"))
            canvas.drawRightString(A4[0]-55, 28, str(doc.page))
        SimpleDocTemplate(output, pagesize=A4, rightMargin=55, leftMargin=55, topMargin=48, bottomMargin=46,
            title=title, author="FinSight").build(story, onFirstPage=footer, onLaterPages=footer)
        return output.getvalue(), "application/pdf"
    if format == "docx":
        from docx import Document
        from docx.shared import Inches, Pt, RGBColor
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        document = Document()
        from docx.shared import Cm
        for section in document.sections:
            section.page_width, section.page_height = Cm(21), Cm(29.7)
            section.top_margin = section.bottom_margin = Cm(1.8)
            section.left_margin = section.right_margin = Cm(2)
        for style in document.styles:
            for border in list(style.element.iter(qn("w:pBdr"))):
                border.getparent().remove(border)
        for name in ("Normal", "Title", "Heading 1", "Heading 2", "Heading 3"):
            style = document.styles[name]
            style.font.name = "Microsoft YaHei"
            style.element.get_or_add_rPr().rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
            style.font.color.rgb = RGBColor(0, 0, 0)
        document.styles["Normal"].font.size = Pt(10.5)
        document.styles["Normal"].paragraph_format.space_after = Pt(8)
        document.styles["Normal"].paragraph_format.line_spacing = 1.35
        document.styles["Title"].font.size = Pt(22)
        document.add_paragraph(title, "Title")
        document.add_paragraph(review_status)
        for kind, value in blocks:
            if kind == "table":
                width = max(len(row) for row in value)
                table = document.add_table(rows=0, cols=width)
                table.style = "Table Grid"
                borders = OxmlElement("w:tblBorders")
                for side in ("top", "left", "bottom", "right", "insideH", "insideV"):
                    border = OxmlElement("w:" + side)
                    for key, attribute in (("val", "single"), ("sz", "4"), ("color", "D9D9D9")):
                        border.set(qn("w:" + key), attribute)
                    borders.append(border)
                table._tbl.tblPr.append(borders)
                for index, row in enumerate(value):
                    cells = table.add_row().cells
                    for column, text in enumerate(row):
                        cells[column].text = text
                    if index == 0:
                        repeat = OxmlElement("w:tblHeader")
                        table.rows[0]._tr.get_or_add_trPr().append(repeat)
                        for cell in cells:
                            shade = OxmlElement("w:shd")
                            shade.set(qn("w:fill"), "E9EFF4")
                            cell._tc.get_or_add_tcPr().append(shade)
                document.add_paragraph()
            else:
                document.add_paragraph(value, "Heading " + str(min(3, int(kind[1]))) if kind.startswith("h") else None)
        if charts:
            document.add_heading("图表与出处", 1)
        for chart in charts:
            document.add_picture(io.BytesIO(chart_png(chart)), width=Inches(6.1))
            document.add_paragraph(chart["interpretation"])
        document.add_heading("来源", 1)
        for ref in references:
            paragraph = document.add_paragraph(ref)
            paragraph.paragraph_format.keep_together = True
            paragraph.paragraph_format.line_spacing = 1.1
            paragraph.paragraph_format.space_after = Pt(5)
            for run in paragraph.runs:
                run.font.size = Pt(8.5)
        document.core_properties.author = "FinSight"
        document.core_properties.title = title
        output = io.BytesIO()
        document.save(output)
        return output.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if format == "pptx":
        # The application must run without Codex's private artifact runtime.
        # python-pptx provides portable, editable Office charts and text.
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        from pptx.enum.shapes import MSO_SHAPE
        presentation = Presentation()
        presentation.slide_width, presentation.slide_height = Inches(13.333), Inches(7.5)
        def slide(title_text, body="", notes=""):
            page = presentation.slides.add_slide(presentation.slide_layouts[6])
            accent = page.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, Inches(.065))
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(24,123,128)
            accent.line.fill.background()
            def textbox(x,y,w,h,content,size,color=(25,45,66)):
                shape = page.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
                frame = shape.text_frame
                frame.word_wrap = True
                for index, line in enumerate(content.split("\n")):
                    p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                    p.text = line
                    p.font.name, p.font.size, p.font.color.rgb = "Microsoft YaHei", Pt(size), RGBColor(*color)
                    p.space_after = Pt(10)
                return shape
            textbox(.7,.45,11.9,1.1,title_text,30)
            if body:
                textbox(.8,1.8,11.7,4.7,body,20)
            textbox(.8,7,11.7,.3,str(len(presentation.slides)),10,(100,115,130))
            textbox(9.3,7,3.2,.3,"FINSIGHT  /  RESEARCH",10,(24,123,128))
            page.notes_slide.notes_text_frame.text = notes
            return page
        slide(title, review_status, "\n\n".join(references))
        for chart in charts:
            page = slide(chart["title"], notes=chart["interpretation"] + "\n\n" + "\n\n".join(references))
            data = CategoryChartData()
            labels = list(dict.fromkeys(p["label"] for p in chart["points"]))
            data.categories = labels
            for series in dict.fromkeys(p["series"] for p in chart["points"]):
                lookup = {p["label"]: p["value"] for p in chart["points"] if p["series"] == series}
                data.add_series(series or chart["unit"], [lookup.get(label) for label in labels])
            figure = page.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS if chart["kind"] == "line" else XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(.8), Inches(1.6), Inches(11.7), Inches(4.7), data).chart
            figure.has_legend = True
            figure.has_title = False
            figure.legend.position = XL_LEGEND_POSITION.BOTTOM
            figure.legend.font.name, figure.legend.font.size = "Microsoft YaHei", Pt(13)
            figure.value_axis.has_title = True
            figure.value_axis.axis_title.text_frame.text = chart["unit"]
            figure.value_axis.tick_labels.number_format = "#,##0.##"
            if chart["kind"] == "bar":
                values = [p["value"] for p in chart["points"]]
                figure.value_axis.minimum_scale = min(0, min(values) * 1.15)
                figure.value_axis.maximum_scale = max(0, max(values) * 1.2) or 1
                for index, series in enumerate(figure.series):
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = RGBColor(*[(24,123,128), (28,54,91), (194,125,62)][index % 3])
                figure.plots[0].has_data_labels = True
                figure.plots[0].data_labels.font.size = Pt(13)
                # Display precision only; preserve the source-backed workbook values.
                figure.plots[0].data_labels.number_format = "#,##0.##"
        current_title, chunks = "研究结论", []
        def flush():
            if not chunks:
                return
            body = "\n\n".join(chunks)
            # No model summarization, truncation or hidden omissions: paginate.
            parts = re.findall(r"[\s\S]{1,320}(?:[。！？\n]|$)|[\s\S]{1,320}", body)
            for part_index, part in enumerate(parts):
                slide(current_title + (f"（续 {part_index}）" if part_index else ""), part.strip(), "\n\n".join(references))
            chunks.clear()
        for kind, value in blocks:
            if kind.startswith("h"):
                flush()
                current_title = value
            elif kind == "table":
                flush()
                width = max(len(row) for row in value)
                for start in range(1, len(value), 5):
                    rows = [value[0], *value[start:start+5]]
                    page = slide(current_title + (f"（表续 {(start-1)//5}）" if start > 1 else ""), notes="\n\n".join(references))
                    table = page.shapes.add_table(len(rows), width, Inches(.8), Inches(1.8), Inches(11.7), Inches(min(4.7, .65 * len(rows)))).table
                    for i, row in enumerate(rows):
                        for j, cell_text in enumerate(row):
                            cell = table.cell(i, j)
                            cell.text = cell_text
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(*( (233,239,244) if i == 0 else (255,255,255) ))
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.name = "Microsoft YaHei"
                                paragraph.font.size = Pt(16)
                                paragraph.font.color.rgb = RGBColor(25,45,66)
            else:
                chunks.append(value)
        flush()
        output = io.BytesIO()
        presentation.save(output)
        return output.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    raise ValueError("unsupported_export_format")
